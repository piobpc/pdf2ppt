import streamlit as st
import pdfplumber
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import openai
from openai import OpenAI
import os
import json
from dotenv import load_dotenv
import time

load_dotenv()

openai.api_key = os.getenv("OPENAI_API_KEY")
client = OpenAI()

st.title("AI-Powered Files Converter || Inteligentny Zamieniacz Plików")

uploaded_file = st.file_uploader("Please upload your PDF presentation | Dodaj prezentację w formacie PDF", type="pdf")

if uploaded_file:
    with open("input.pdf", "wb") as f:
        f.write(uploaded_file.read())

    st.write("Converting... | W trakcie zamiany...")

    try:
        prs = Presentation()
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        image_counter = 0
        image_files = []

        with pdfplumber.open("input.pdf") as pdf:
            for page_number, page in enumerate(pdf.pages, start=1):
                page_text = page.extract_text()

                if not page_text:
                    continue  # Skip empty pages

                layout_elements = {
                    "TextBlocks": page.extract_words(extra_attrs=["fontname", "size", "x0", "x1", "top", "bottom"]),
                    "Images": [],
                    "Rects": page.rects,
                    "Lines": page.lines
                }

                # Extract and save images
                for img in page.images:
                    bbox = (img['x0'], img['top'], img['x1'], img['bottom'])
                    cropped_image = page.crop(bbox).to_image(resolution=150)
                    image_path = f"image_{image_counter}.png"
                    cropped_image.save(image_path)
                    layout_elements["Images"].append({"path": image_path, "x0": img['x0'], "top": img['top'], "x1": img['x1'], "bottom": img['bottom']})
                    image_files.append(image_path)
                    image_counter += 1

                page_text_trimmed = page_text[:3000]

                prompt = f"""
Poniżej znajduje się layout slajdu PDF (teksty, obrazy, prostokąty, linie) oraz wyodrębniony tekst.

Twoje zadanie:
1. Odtwórz strukturę slajdu jako edytowalny PowerPoint.
2. Określ pozycje elementów (x, y, szerokość, wysokość) w skali 0-1 (gdzie 1 to szerokość/wysokość slajdu).
3. Określ typy elementów: Tytuł, BulletPoint, Obraz, Prostokąt, Linia, Wykres, Tabela.
4. Jeśli widzisz wykres lub tabelę, podaj dane przykładowe.

Zwróć JSON w formacie:
{{
    "Elements": [
        {{
            "Type": "Title" / "BulletPoint" / "Image" / "Shape" / "Line" / "Chart" / "Table",
            "Text": "...", 
            "x": 0.1,
            "y": 0.2,
            "width": 0.8,
            "height": 0.1,
            "ImageIndex": 0,  # Jeśli Type == Image
            "Style": {{"FontSize": 32, "Bold": true, "Color": "#000000"}}  
        }},
        ...
    ]
}}

Layout Elements:
{json.dumps(layout_elements, indent=2)}
"""

                layout_data = None
                for attempt in range(3):
                    response = client.chat.completions.create(
                        model="gpt-4o",
                        messages=[{"role": "user", "content": prompt}]
                    )
                    content = response.choices[0].message.content.strip()

                    if content.startswith("{"):
                        try:
                            layout_data = json.loads(content)
                            break
                        except json.JSONDecodeError:
                            pass
                    time.sleep(1)

                if not layout_data:
                    st.warning(f"AI nie zwróciło poprawnego JSON dla strony {page_number}. Tworzę awaryjny slajd z tekstem.")
                    layout_data = {
                        "Elements": [
                            {
                                "Type": "Title",
                                "Text": f"Page {page_number}",
                                "x": 0.05,
                                "y": 0.05,
                                "width": 0.9,
                                "height": 0.1,
                                "Style": {"FontSize": 32, "Bold": True, "Color": "#000000"}
                            }
                        ] + [
                            {
                                "Type": "BulletPoint",
                                "Text": line.strip(),
                                "x": 0.1,
                                "y": 0.2 + idx * 0.05,
                                "width": 0.8,
                                "height": 0.05,
                                "Style": {"FontSize": 18, "Bold": False, "Color": "#000000"}
                            } for idx, line in enumerate(page_text_trimmed.split('\n')[:10])
                        ]
                    }

                # --- Build slide ---
                slide_layout = prs.slide_layouts[6]  # Blank slide
                slide = prs.slides.add_slide(slide_layout)

                for element in layout_data.get("Elements", []):
                    el_type = element.get("Type")
                    text = element.get("Text", "")
                    x = slide_width * element.get("x", 0)
                    y = slide_height * element.get("y", 0)
                    width = slide_width * element.get("width", 1)
                    height = slide_height * element.get("height", 1)
                    style = element.get("Style", {})

                    if el_type in ["Title", "BulletPoint"]:
                        textbox = slide.shapes.add_textbox(x, y, width, height)
                        text_frame = textbox.text_frame
                        p = text_frame.paragraphs[0]
                        run = p.add_run()
                        run.text = text
                        if "FontSize" in style:
                            run.font.size = Pt(style["FontSize"])
                        if style.get("Bold"):
                            run.font.bold = True
                        if "Color" in style:
                            hex_color = style["Color"].lstrip('#')
                            rgb = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
                            run.font.color.rgb = RGBColor(*rgb)

                    elif el_type == "Image":
                        image_index = element.get("ImageIndex", 0)
                        if 0 <= image_index < len(image_files):
                            slide.shapes.add_picture(image_files[image_index], x, y, width, height)

                    elif el_type == "Shape":
                        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, width, height)
                        shape.fill.solid()
                        shape.fill.fore_color.rgb = RGBColor(200, 200, 200)

                    elif el_type == "Line":
                        line = slide.shapes.add_shape(MSO_SHAPE.LINE_INVERSE, x, y, width, height)
                        line.line.color.rgb = RGBColor(0, 0, 0)

        prs.save("output.pptx")

        with open("output.pptx", "rb") as f:
            st.download_button("Download | Pobierz", f, file_name="converted.pptx")

    except Exception as e:
        st.error(f"Error | Błąd: {e}")

    # Cleanup temporary files
    if os.path.exists("input.pdf"):
        os.remove("input.pdf")
    if os.path.exists("output.pptx"):
        os.remove("output.pptx")
    for img_file in image_files:
        if os.path.exists(img_file):
            os.remove(img_file)
